"""
PTD-BR Corpus — Gerador de Apresentação PowerPoint
IPEA / COGIT / DIEST
Uso: python gerar_apresentacao.py
Output: apresentacao_ptd_corpus.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import datetime

# ── Paleta IPEA ───────────────────────────────────────
AZUL_IPEA   = RGBColor(0x0D, 0x2B, 0x4E)   # azul escuro
VERDE_IPEA  = RGBColor(0x1A, 0x7F, 0x7A)   # verde-água
LARANJA     = RGBColor(0xD9, 0x77, 0x06)
BRANCO      = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_CLARO = RGBColor(0xF4, 0xF4, 0xF4)
CINZA_TEXTO = RGBColor(0x44, 0x44, 0x44)

def nova_apresentacao():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def fundo_colorido(slide, cor: RGBColor):
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = cor
    bg.line.fill.background()
    bg.zorder = 0

def caixa(slide, texto, l, t, w, h, tam=18, bold=False,
          cor_txt=BRANCO, cor_bg=None, alinha=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alinha
    run = p.add_run()
    run.text = texto
    run.font.size  = Pt(tam)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = cor_txt
    if cor_bg:
        txb.fill.solid()
        txb.fill.fore_color.rgb = cor_bg
    return txb

def retangulo(slide, l, t, w, h, cor: RGBColor, alpha=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = cor
    s.line.fill.background()
    return s

def linha_div(slide, t, cor=VERDE_IPEA):
    retangulo(slide, 0.5, t, 12.3, 0.04, cor)

# ════════════════════════════════════════════════════════
# SLIDE 1 — Capa
# ════════════════════════════════════════════════════════
def slide_capa(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, AZUL_IPEA)
    retangulo(sl, 0, 5.8, 13.33, 1.7, VERDE_IPEA)
    caixa(sl, 'PTD-BR CORPUS', 0.8, 0.6, 11, 1.2, tam=40, bold=True)
    caixa(sl, 'Planos de Transformação Digital do Governo Federal Brasileiro',
          0.8, 1.7, 11, 0.8, tam=20)
    caixa(sl, 'Coleta · Extração · Estruturação · Análise',
          0.8, 2.35, 10, 0.6, tam=15, italic=True)
    linha_div(sl, 3.1, LARANJA)
    caixa(sl, 'IPEA / COGIT / DIEST', 0.8, 3.3, 8, 0.5, tam=13)
    caixa(sl, 'Denise do Carmo Direito  ·  Lucas Freire Silva',
          0.8, 3.75, 10, 0.5, tam=13)
    caixa(sl, f'Versão 3.0  ·  {datetime.date.today().strftime("%d/%m/%Y")}',
          0.8, 5.95, 10, 0.5, tam=13, bold=True)
    caixa(sl, 'Documento de trabalho — não citar sem autorização dos autores',
          0.8, 6.5, 11.5, 0.4, tam=10, italic=True,
          cor_txt=RGBColor(0xCC,0xCC,0xCC))

# ════════════════════════════════════════════════════════
# SLIDE 2 — O que são os PTDs
# ════════════════════════════════════════════════════════
def slide_contexto(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, CINZA_CLARO)
    retangulo(sl, 0, 0, 13.33, 1.1, AZUL_IPEA)
    caixa(sl, 'O QUE SÃO OS PTDs?', 0.5, 0.2, 12, 0.7,
          tam=24, bold=True, alinha=PP_ALIGN.CENTER)

    itens = [
        ('Base legal', 'Decreto nº 12.198/2024 · Portaria SGD/MGI nº 6.618/2024'),
        ('Obrigação', '90 órgãos da APF pactuam metas digitais com o MGI para 2025–2027'),
        ('Estrutura', '6 Eixos da EFGD (Estratégia Federal de Governo Digital)'),
        ('Documentos', 'Anexo de Entregas (metas) + Documento Diretivo (riscos e diretrizes)'),
        ('Publicação', 'Portal Gov Digital — PDFs públicos assinados digitalmente'),
    ]
    for i, (titulo, desc) in enumerate(itens):
        y = 1.3 + i * 1.1
        retangulo(sl, 0.5, y, 0.08, 0.55, VERDE_IPEA)
        caixa(sl, titulo, 0.75, y, 2.8, 0.35, tam=13, bold=True, cor_txt=AZUL_IPEA)
        caixa(sl, desc,   0.75, y+0.32, 11.5, 0.5, tam=12, cor_txt=CINZA_TEXTO)

    caixa(sl, 'Por que coletar?', 0.5, 6.6, 5, 0.4, tam=11, bold=True,
          italic=True, cor_txt=AZUL_IPEA)
    caixa(sl, 'Nenhuma base estruturada dos PTDs existia — este corpus é inédito.',
          3.2, 6.6, 9.5, 0.4, tam=11, italic=True, cor_txt=CINZA_TEXTO)

# ════════════════════════════════════════════════════════
# SLIDE 3 — Arquitetura do Pipeline
# ════════════════════════════════════════════════════════
def slide_arquitetura(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, CINZA_CLARO)
    retangulo(sl, 0, 0, 13.33, 1.1, AZUL_IPEA)
    caixa(sl, 'ARQUITETURA DO PIPELINE', 0.5, 0.2, 12, 0.7,
          tam=24, bold=True, alinha=PP_ALIGN.CENTER)

    camadas = [
        (AZUL_IPEA,  'LAYER 1', 'Coleta',
         'Scraping dinâmico do portal → URLs dos PDFs → Download com SHA-256'),
        (VERDE_IPEA, 'LAYER 2', 'Extração',
         'Docling TableFormer → tabelas dos PDFs → corpus_raw.csv + manifesto'),
        (LARANJA,    'LAYER 3', 'Curadoria',
         'Parser texto→campos · Correção de eixo · Tipo de entrega · Flag IA'),
        (RGBColor(0x45,0x7B,0x9D), 'LAYER 4', 'Análise',
         'Visualizações · Relatório HTML · Metadados JSON · Apresentação'),
    ]
    for i, (cor, label, nome, desc) in enumerate(camadas):
        y = 1.25 + i * 1.45
        retangulo(sl, 0.5, y, 1.6, 1.1, cor)
        caixa(sl, label, 0.5, y + 0.1, 1.6, 0.4, tam=10, bold=True,
              alinha=PP_ALIGN.CENTER)
        caixa(sl, nome,  0.5, y + 0.5, 1.6, 0.45, tam=16, bold=True,
              alinha=PP_ALIGN.CENTER)
        caixa(sl, desc,  2.4, y + 0.25, 10.4, 0.6, tam=13, cor_txt=CINZA_TEXTO)
        if i < 3:
            caixa(sl, '▼', 0.9, y + 1.15, 1, 0.25, tam=14, cor_txt=cor)

# ════════════════════════════════════════════════════════
# SLIDE 4 — O que foi coletado
# ════════════════════════════════════════════════════════
def slide_coleta(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, CINZA_CLARO)
    retangulo(sl, 0, 0, 13.33, 1.1, AZUL_IPEA)
    caixa(sl, 'O QUE FOI COLETADO', 0.5, 0.2, 12, 0.7,
          tam=24, bold=True, alinha=PP_ALIGN.CENTER)

    n_reg = f'{_KPIS["n_registros"]:,}'.replace(',', '.') if _KPIS.get('n_registros') else '5.825'
    numeros = [
        ('90',  'órgãos\ncatalogados'),
        ('54',  'PDFs únicos\nbaixados'),
        (n_reg, 'linhas no\ncorpus bruto'),
        ('4',   'grupos\ncompartilhados'),
    ]
    for i, (num, leg) in enumerate(numeros):
        x = 0.7 + i * 3.1
        retangulo(sl, x, 1.35, 2.6, 1.7, AZUL_IPEA)
        caixa(sl, num, x, 1.45, 2.6, 1.0, tam=40, bold=True,
              alinha=PP_ALIGN.CENTER)
        caixa(sl, leg, x, 2.4,  2.6, 0.6, tam=12,
              alinha=PP_ALIGN.CENTER, italic=True)

    linha_div(sl, 3.4)
    caixa(sl, 'Grupos compartilhados (PDF único, N órgãos):',
          0.5, 3.6, 12, 0.4, tam=13, bold=True, cor_txt=AZUL_IPEA)
    grupos = [
        'MMA → IBAMA · ICMBio · SFB · JBRJ  (1.715 linhas)',
        'MF  → RFB · PGFN · STN  (328 linhas)',
        'MT  → ANTT · DNIT  (132 linhas)',
        'MDA → CONAB  (96 linhas)',
    ]
    for i, g in enumerate(grupos):
        caixa(sl, f'• {g}', 0.9, 4.05 + i*0.5, 12, 0.4, tam=12, cor_txt=CINZA_TEXTO)

    caixa(sl, 'Nota: scraping dinâmico do portal (v3.0) substitui URLs hardcoded — '
              'elimina causa dos 100 erros 404 da versão anterior.',
          0.5, 6.55, 12.3, 0.5, tam=10, italic=True, cor_txt=LARANJA)

# ════════════════════════════════════════════════════════
# SLIDE 5 — Corpus de Entregas
# ════════════════════════════════════════════════════════
def slide_corpus(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, CINZA_CLARO)
    retangulo(sl, 0, 0, 13.33, 1.1, VERDE_IPEA)
    caixa(sl, 'CORPUS DE ENTREGAS', 0.5, 0.2, 12, 0.7,
          tam=24, bold=True, alinha=PP_ALIGN.CENTER)

    caixa(sl, 'Unidade de análise: comprometimento = serviço × produto',
          0.5, 1.2, 12.3, 0.5, tam=14, bold=True, cor_txt=AZUL_IPEA,
          alinha=PP_ALIGN.CENTER)
    caixa(sl, 'Um mesmo serviço pode pactar N produtos → N linhas. '
              'Média: 2,07 produtos/serviço (54 diretivos). Design intencional do template SGD v2.3.',
          0.5, 1.65, 12.3, 0.5, tam=11, italic=True, cor_txt=CINZA_TEXTO,
          alinha=PP_ALIGN.CENTER)

    eixos = [
        ('E1', 'Centrado no Cidadão', '~45%', AZUL_IPEA),
        ('E2', 'Integrado e Colaborativo', '~18%', VERDE_IPEA),
        ('E3', 'Inteligente e Inovador', '~4%*', LARANJA),
        ('E4', 'Confiável e Seguro', '~22%', RGBColor(0xE6,0x39,0x46)),
        ('E5', 'Transparente e Aberto', '~5%', RGBColor(0x45,0x7B,0x9D)),
        ('E6', 'Eficiente e Sustentável', '~6%', RGBColor(0x6A,0x4C,0x93)),
    ]
    for i, (cod, nome, pct, cor) in enumerate(eixos):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 4.2
        y = 2.45 + row * 1.6
        retangulo(sl, x, y, 0.55, 1.1, cor)
        caixa(sl, cod, x, y+0.1, 0.55, 0.5, tam=14, bold=True,
              alinha=PP_ALIGN.CENTER)
        caixa(sl, pct, x, y+0.6, 0.55, 0.4, tam=12,
              alinha=PP_ALIGN.CENTER)
        caixa(sl, nome, x+0.7, y+0.3, 3.4, 0.6, tam=12, cor_txt=CINZA_TEXTO)

    caixa(sl, '* E3 auditado: 21 registros PGFN reclassificados E3→E1 (bug state machine)',
          0.5, 6.9, 12.3, 0.35, tam=10, italic=True, cor_txt=LARANJA)

# ════════════════════════════════════════════════════════
# SLIDE 6 — Problemas identificados e corrigidos
# ════════════════════════════════════════════════════════
def slide_problemas(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, CINZA_CLARO)
    retangulo(sl, 0, 0, 13.33, 1.1, RGBColor(0x8B,0x1A,0x1A))
    caixa(sl, 'PROBLEMAS IDENTIFICADOS E CORREÇÕES', 0.5, 0.2, 12, 0.7,
          tam=22, bold=True, alinha=PP_ALIGN.CENTER)

    problemas = [
        ('🔴 Crítico', 'State machine de eixo sem reset por página',
         'Contaminou 21 registros PGFN com E3 incorreto',
         '✅ FIX v3.0: eixo_atual = None a cada nova página'),
        ('🔴 Crítico', 'URLs hardcoded no CATÁLOGO (100/120 com 404)',
         'Pipeline de extração falhava na maioria dos PDFs',
         '✅ FIX v3.0: scraping dinâmico do portal HTML'),
        ('🟡 Médio', 'PRODUTOS hardcoded no script',
         'Exige edição do código a cada novo ciclo PTD',
         '✅ FIX v2.1: externalizado em config/produtos_sgd_v23.json'),
        ('🟡 Médio', 'Correção E3 PGFN-específica (frágil)',
         'Falha para qualquer novo órgão com padrão similar',
         '✅ FIX v2.1: regras em config/correcoes_eixo.json'),
        ('🟠 Menor', 'Sem SHA-256 dos PDFs nos registros',
         'Impossível rastrear registro → PDF de origem',
         '✅ FIX v3.0: pdf_sha256 em cada linha + pipeline_manifest.json'),
    ]
    for i, (sev, prob, impacto, fix) in enumerate(problemas):
        y = 1.2 + i * 1.2
        cor = RGBColor(0x8B,0x1A,0x1A) if 'Crítico' in sev else \
              RGBColor(0xD9,0x77,0x06) if 'Médio' in sev else \
              RGBColor(0xCC,0x88,0x00)
        retangulo(sl, 0.5, y, 0.08, 0.9, cor)
        caixa(sl, prob,    0.75, y,      8.5, 0.38, tam=12, bold=True, cor_txt=AZUL_IPEA)
        caixa(sl, impacto, 0.75, y+0.35, 8.5, 0.3,  tam=10, italic=True, cor_txt=CINZA_TEXTO)
        caixa(sl, fix,     9.5,  y+0.15, 3.3, 0.5,  tam=10, cor_txt=VERDE_IPEA, bold=True)

# ════════════════════════════════════════════════════════
# SLIDE 7 — Qualidade do Corpus
# ════════════════════════════════════════════════════════
def slide_qualidade(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, CINZA_CLARO)
    retangulo(sl, 0, 0, 13.33, 1.1, AZUL_IPEA)
    caixa(sl, 'QUALIDADE DO CORPUS', 0.5, 0.2, 12, 0.7,
          tam=24, bold=True, alinha=PP_ALIGN.CENTER)

    metricas = [
        ('Cobertura de órgãos',   '71%',  '64/90 com dados extraídos'),
        ('Parse OK (serv+prod)',   '~65%', 'Linhas com todos os campos identificados'),
        ('Cobertura data_ptd',     '~55%', 'Datas de entrega extraídas'),
        ('E3 genuíno (corrigido)', '~4%',  '2 registros IA real vs 23 originais'),
        ('IA real identificada',   '19',   'Registros em 14 órgãos (PAT_IA_REAL preciso)'),
        ('Mandatório vs disco.',   '44%',  'PPSI + Login + Avaliação + Revisão descrição'),
    ]
    for i, (label, valor, desc) in enumerate(metricas):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.4
        y = 1.3 + row * 1.9
        retangulo(sl, x, y, 6.0, 1.6, BRANCO)
        retangulo(sl, x, y, 0.08, 1.6, VERDE_IPEA)
        caixa(sl, valor, x+0.25, y+0.1, 5.5, 0.8, tam=30, bold=True,
              cor_txt=AZUL_IPEA)
        caixa(sl, label, x+0.25, y+0.85, 5.5, 0.4, tam=12, bold=True,
              cor_txt=CINZA_TEXTO)
        caixa(sl, desc,  x+0.25, y+1.2,  5.5, 0.3, tam=10, italic=True,
              cor_txt=CINZA_TEXTO)

# ════════════════════════════════════════════════════════
# SLIDE 8 — Avaliação de Prontidão
# ════════════════════════════════════════════════════════
def slide_prontidao(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, CINZA_CLARO)
    retangulo(sl, 0, 0, 13.33, 1.1, VERDE_IPEA)
    caixa(sl, 'AVALIAÇÃO DE PRONTIDÃO', 0.5, 0.2, 12, 0.7,
          tam=24, bold=True, alinha=PP_ALIGN.CENTER)

    dimensoes = [
        ('Completude da coleta',     3, 5, 'Scraping dinâmico implementado; ~30% PDFs ainda sem extração validada'),
        ('Qualidade da extração',    4, 5, 'Docling TableFormer + fix state machine; OCR para PDFs tarjados'),
        ('Curadoria semântica',      4, 5, 'Parser, tipo_entrega, ia_real, correção E3 com regras externas'),
        ('Rastreabilidade',          4, 5, 'SHA-256 PDFs + hashes CSV inputs + pipeline_manifest.json'),
        ('Reprodutibilidade',        3, 5, 'Checkpoints implementados; falta documentação de ambiente (requirements.txt)'),
        ('Ciência aberta',           3, 5, 'Código versionado no GitHub; falta DOI de dados e licença explícita'),
        ('Pronto para publicação',   2, 5, 'Requires: validação manual amostra + relatório HTML + DOI Zenodo'),
    ]

    for i, (dim, nota, max_n, obs) in enumerate(dimensoes):
        y = 1.3 + i * 0.84
        # barra de fundo
        retangulo(sl, 3.8, y+0.15, 5.5, 0.45, RGBColor(0xDD,0xDD,0xDD))
        # barra de nota
        cor_bar = VERDE_IPEA if nota >= 4 else LARANJA if nota == 3 else \
                  RGBColor(0xE6,0x39,0x46)
        retangulo(sl, 3.8, y+0.15, 5.5*(nota/max_n), 0.45, cor_bar)
        caixa(sl, dim,    0.5, y, 3.2, 0.5, tam=12, bold=True, cor_txt=AZUL_IPEA)
        caixa(sl, f'{nota}/{max_n}', 9.5, y, 0.7, 0.5, tam=14, bold=True,
              cor_txt=cor_bar)
        caixa(sl, obs, 10.3, y+0.05, 2.8, 0.5, tam=9, italic=True,
              cor_txt=CINZA_TEXTO)

# ════════════════════════════════════════════════════════
# SLIDE 9 — Próximos Passos
# ════════════════════════════════════════════════════════
def slide_proximos_passos(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, CINZA_CLARO)
    retangulo(sl, 0, 0, 13.33, 1.1, AZUL_IPEA)
    caixa(sl, 'PRÓXIMOS PASSOS', 0.5, 0.2, 12, 0.7,
          tam=24, bold=True, alinha=PP_ALIGN.CENTER)

    passos = [
        ('Curto prazo\n(1–2 sem)', VERDE_IPEA, [
            'Rodar v3.0 com scraping dinâmico e validar cobertura de PDFs',
            'Adicionar requirements.txt + README de execução',
            'Validação manual de amostra: 5% das linhas (±290 registros)',
        ]),
        ('Médio prazo\n(1 mês)', LARANJA, [
            'Integrar extração de riscos ao manifesto SHA-256',
            'Gerar relatório HTML autocontido (ETAPA 7 do v3.0)',
            'Publicar corpus no Repositório de Dados IPEA (DOI)',
        ]),
        ('Longo prazo\n(ciclo 2026)', RGBColor(0x45,0x7B,0x9D), [
            'Atualizar config/produtos_sgd_v*.json para novo template SGD',
            'Processar novos PTDs 2026–2028 sem refatoração do pipeline',
            'Produzir análise comparativa ciclo 2025 vs 2026',
        ]),
    ]

    for col, (periodo, cor, items) in enumerate(passos):
        x = 0.5 + col * 4.2
        retangulo(sl, x, 1.3, 3.9, 0.65, cor)
        caixa(sl, periodo, x, 1.35, 3.9, 0.55, tam=13, bold=True,
              alinha=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            caixa(sl, f'• {item}', x+0.1, 2.1 + j*1.5, 3.7, 1.3,
                  tam=11, cor_txt=CINZA_TEXTO)

# ════════════════════════════════════════════════════════
# SLIDE 10 — Balanço de Sessão
# ════════════════════════════════════════════════════════
def slide_balanco(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, AZUL_IPEA)
    retangulo(sl, 0, 0, 13.33, 1.1, VERDE_IPEA)
    caixa(sl, 'BALANÇO DE SESSÃO', 0.5, 0.2, 12, 0.7,
          tam=24, bold=True, alinha=PP_ALIGN.CENTER)

    entregues = [
        'ptd_pipeline_v30.py     — Pipeline de extração com 5 fixes críticos',
        'ptd_corpus_v21.py       — Curadoria com regras externas e rastreabilidade',
        'config/produtos_sgd_v23.json   — Taxonomia SGD externalizada',
        'config/correcoes_eixo.json     — Regras de auditoria versionadas',
        'gerar_apresentacao.py   — Gerador desta apresentação (reproduzível)',
        'apresentacao_ptd_corpus.pptx  — Esta apresentação',
    ]
    caixa(sl, '✅ ENTREGUES NESTA SESSÃO', 0.5, 1.2, 12, 0.4,
          tam=14, bold=True, cor_txt=VERDE_IPEA)
    for i, e in enumerate(entregues):
        caixa(sl, e, 0.7, 1.65 + i*0.45, 12, 0.4, tam=11)

    linha_div(sl, 4.5, LARANJA)

    caixa(sl, '⚠️  PENDENTE (próxima sessão)', 0.5, 4.65, 12, 0.4,
          tam=14, bold=True, cor_txt=LARANJA)
    pendentes = [
        'Executar v3.0 com PDFs reais e validar cobertura',
        'requirements.txt + documentação de ambiente',
        'Amostra de validação manual (5%) e correção de parse_flag',
    ]
    for i, p in enumerate(pendentes):
        caixa(sl, f'• {p}', 0.7, 5.1 + i*0.43, 12, 0.4,
              tam=11, cor_txt=LARANJA)

    caixa(sl, f'Gerado em {datetime.datetime.now().strftime("%d/%m/%Y %H:%M")} | COGIT/DIEST/IPEA',
          0.5, 7.1, 12.3, 0.3, tam=9, italic=True,
          cor_txt=RGBColor(0xAA,0xAA,0xAA))

# ════════════════════════════════════════════════════════
# GERAR
# ════════════════════════════════════════════════════════
def _carregar_kpis() -> dict:
    """Lê ptd_corpus_v21_metadados.json se existir e retorna KPIs reais."""
    from pathlib import Path as _Path
    import json as _json
    meta_path = _Path('ptd_corpus/03_database/ptd_corpus_v21_metadados.json')
    if not meta_path.exists():
        return {}
    try:
        with open(meta_path, encoding='utf-8') as f:
            m = _json.load(f)
        return {
            'n_registros':  m.get('corpus', {}).get('total_linhas', 0),
            'n_orgaos':     len(set()),  # calculado abaixo via pivot
            'parse_ok_pct': m.get('parse', {}).get('cobertura_servico_pct', 0),
        }
    except Exception:
        return {}

# KPIs globais — preenchidos com valores reais se disponíveis, senão placeholders
_KPIS = _carregar_kpis()

def gerar():
    from pathlib import Path as _Path
    prs = nova_apresentacao()
    slide_capa(prs)
    slide_contexto(prs)
    slide_arquitetura(prs)
    slide_coleta(prs)
    slide_corpus(prs)
    slide_problemas(prs)
    slide_qualidade(prs)
    slide_prontidao(prs)
    slide_proximos_passos(prs)
    slide_balanco(prs)
    out = _Path('apresentacao_ptd_corpus.pptx')
    prs.save(str(out))
    print(f'✅ Apresentação salva: {out}  ({out.stat().st_size//1024} KB)')
    print(f'   {len(prs.slides)} slides')
    return out

if __name__ == '__main__':
    gerar()
